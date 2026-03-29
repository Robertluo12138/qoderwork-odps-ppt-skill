from __future__ import annotations

import argparse
import json
import platform
from pathlib import Path
from typing import Any

from common import deep_render, load_yaml, output_dir_from_config


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Render an editable PPTX from a report payload and slide plan."
    )
    parser.add_argument("--config", required=True, help="Path to the YAML report config.")
    parser.add_argument("--payload", required=True, help="Path to report_payload.json.")
    parser.add_argument("--plan", required=True, help="Path to slide_plan.generated.json.")
    parser.add_argument("--output", default="", help="Optional explicit output path.")
    return parser.parse_args()


def load_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def hex_to_rgb(color: str) -> tuple[int, int, int]:
    clean = color.strip().lstrip("#")
    if len(clean) != 6:
        raise ValueError(f"Expected 6-digit hex color, got: {color}")
    return tuple(int(clean[index : index + 2], 16) for index in (0, 2, 4))


def to_float(value: Any) -> float | None:
    try:
        text = str(value).replace(",", "").strip()
        if not text:
            return None
        return float(text)
    except ValueError:
        return None


def ensure_dependencies() -> None:
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is not installed. Run scripts/bootstrap_env.py first."
        ) from exc


def resolve_font_family(raw_theme: dict[str, Any]) -> str:
    system_name = platform.system().lower()
    if system_name == "darwin":
        return str(
            raw_theme.get("font_family_mac")
            or raw_theme.get("font_family")
            or "Aptos"
        )
    if system_name == "windows":
        return str(
            raw_theme.get("font_family_windows")
            or raw_theme.get("font_family")
            or "Aptos"
        )
    return str(
        raw_theme.get("font_family_other")
        or raw_theme.get("font_family")
        or "Aptos"
    )


def resolve_cjk_font_family(raw_theme: dict[str, Any]) -> str:
    system_name = platform.system().lower()
    if system_name == "darwin":
        return str(
            raw_theme.get("cjk_font_family_mac")
            or raw_theme.get("cjk_font_family")
            or "Hiragino Sans GB"
        )
    if system_name == "windows":
        return str(
            raw_theme.get("cjk_font_family_windows")
            or raw_theme.get("cjk_font_family")
            or "Microsoft YaHei"
        )
    return str(
        raw_theme.get("cjk_font_family_other")
        or raw_theme.get("cjk_font_family")
        or "Noto Sans CJK SC"
    )


def build_theme(raw_theme: dict[str, Any]) -> dict[str, Any]:
    chart_palette = raw_theme.get("chart_palette") or [
        "FF6200",
        "FF8A3D",
        "FFA020",
        "FFB85C",
        "FFDCC8",
    ]
    return {
        "primary_color": raw_theme.get("primary_color", "FF6200"),
        "accent_color": raw_theme.get("accent_color", "FFA020"),
        "soft_fill_color": raw_theme.get("soft_fill_color", "FFF4EB"),
        "warm_fill_color": raw_theme.get("warm_fill_color", "FFDCC8"),
        "text_color": raw_theme.get("text_color", "1F1F1F"),
        "muted_text_color": raw_theme.get("muted_text_color", "666666"),
        "line_color": raw_theme.get("line_color", "EAEAEA"),
        "card_bg_color": raw_theme.get("card_bg_color", "FAFAFA"),
        "background_color": raw_theme.get("background_color", "FFFFFF"),
        "font_family": resolve_font_family(raw_theme),
        "cjk_font_family": resolve_cjk_font_family(raw_theme),
        "language": str(raw_theme.get("language", "zh-CN")),
        "title_font_size": int(raw_theme.get("title_font_size", 24)),
        "hero_title_font_size": int(raw_theme.get("hero_title_font_size", 28)),
        "subtitle_font_size": int(raw_theme.get("subtitle_font_size", 18)),
        "body_font_size": int(raw_theme.get("body_font_size", 16)),
        "table_header_font_size": int(raw_theme.get("table_header_font_size", 11)),
        "table_body_font_size": int(raw_theme.get("table_body_font_size", 10)),
        "footer_font_size": int(raw_theme.get("footer_font_size", 10)),
        "chart_palette": [str(color).lstrip("#") for color in chart_palette],
        "logo_path": str(raw_theme.get("logo_path", "taobao_flash_sale_logo.png")),
    }


def contains_cjk(text: str) -> bool:
    for char in text:
        codepoint = ord(char)
        if (
            0x3400 <= codepoint <= 0x4DBF
            or 0x4E00 <= codepoint <= 0x9FFF
            or 0xF900 <= codepoint <= 0xFAFF
            or 0x3040 <= codepoint <= 0x30FF
            or 0xAC00 <= codepoint <= 0xD7AF
        ):
            return True
    return False


def set_run_typefaces(
    run,
    latin_font_name: str,
    cjk_font_name: str,
    language: str,
) -> None:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    rPr = run._r.get_or_add_rPr()
    rPr.set("lang", language)

    def upsert_font(tag: str, typeface: str) -> None:
        child = rPr.find(qn(tag))
        if child is None:
            child = OxmlElement(tag)
            rPr.insert_element_before(
                child,
                "a:hlinkClick",
                "a:hlinkMouseOver",
                "a:rtl",
                "a:extLst",
            )
        child.set("typeface", typeface)

    upsert_font("a:latin", latin_font_name)
    upsert_font("a:ea", cjk_font_name or latin_font_name)
    upsert_font("a:cs", cjk_font_name or latin_font_name)


def apply_run_style(
    run,
    font_name: str,
    font_size: int,
    rgb: tuple[int, int, int],
    bold: bool = False,
    cjk_font_name: str = "",
    language: str = "zh-CN",
):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    run_text = getattr(run, "text", "") or ""
    preferred_font = cjk_font_name if cjk_font_name and contains_cjk(run_text) else font_name
    run.font.name = preferred_font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*rgb)
    set_run_typefaces(run, font_name, cjk_font_name or preferred_font, language)


def add_textbox(slide, left: float, top: float, width: float, height: float):
    from pptx.util import Inches

    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def add_title_textbox(slide, title: str, theme: dict[str, Any]):
    box = add_textbox(slide, 0.6, 0.45, 12.1, 0.8)
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    apply_run_style(
        run,
        theme["font_family"],
        theme["title_font_size"],
        hex_to_rgb(theme["text_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )
    return box


def add_bullets_box(
    slide,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    theme: dict[str, Any],
    font_size: int | None = None,
):
    box = add_textbox(slide, left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True

    if not bullets:
        bullets = ["本页还没有填内容。"]

    first = True
    for bullet in bullets:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        paragraph.level = 0
        run = paragraph.add_run()
        run.text = bullet
        apply_run_style(
            run,
            theme["font_family"],
            font_size or theme["body_font_size"],
            hex_to_rgb(theme["text_color"]),
            bold=False,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        first = False
    return box


def add_table_shape(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left: float,
    top: float,
    width: float,
    height: float,
    theme: dict[str, Any],
):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    row_count = max(2, len(rows) + 1)
    col_count = max(1, len(headers))
    table_shape = slide.shapes.add_table(
        row_count,
        col_count,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table
    primary = RGBColor(*hex_to_rgb(theme["primary_color"]))
    text_rgb = RGBColor(*hex_to_rgb(theme["text_color"]))

    for idx, header in enumerate(headers):
        cell = table.cell(0, idx)
        cell.text = header
        paragraph = cell.text_frame.paragraphs[0]
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        if not paragraph.runs:
            run.text = header
        apply_run_style(
            run,
            theme["font_family"],
            theme["table_header_font_size"],
            (255, 255, 255),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary

    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            paragraph = cell.text_frame.paragraphs[0]
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            apply_run_style(
                run,
                theme["font_family"],
                theme["table_body_font_size"],
                hex_to_rgb(theme["text_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )

    return table_shape


def add_footer(slide, text: str, theme: dict[str, Any]):
    logo_path = theme.get("_logo_resolved_path")
    if logo_path and Path(logo_path).exists():
        from pptx.util import Inches

        slide.shapes.add_picture(
            logo_path, Inches(11.6), Inches(7.02), height=Inches(0.22)
        )
    if not text:
        return

    box = add_textbox(slide, 0.6, 7.0, 11.0, 0.3)
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    apply_run_style(
        run,
        theme["font_family"],
        theme["footer_font_size"],
        hex_to_rgb(theme["text_color"]),
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )


def resolve_logo_path(config_path: Path, theme: dict[str, Any]) -> str | None:
    """Resolve the logo path from theme config, checking multiple locations."""
    logo_rel = theme.get("logo_path", "")
    if not logo_rel:
        return None
    candidate = config_path.parent / logo_rel
    if candidate.exists():
        return str(candidate)
    skill_assets = Path(__file__).resolve().parent.parent / "assets"
    candidate = skill_assets / logo_rel
    if candidate.exists():
        return str(candidate)
    return None


def add_logo_to_slide(
    slide, theme: dict[str, Any], left: float, top: float, height: float
) -> None:
    """Add the brand logo to a slide if the resolved logo path exists."""
    logo_path = theme.get("_logo_resolved_path")
    if not logo_path or not Path(logo_path).exists():
        return
    from pptx.util import Inches

    slide.shapes.add_picture(
        logo_path, Inches(left), Inches(top), height=Inches(height)
    )


def normalize_rows(
    row_dicts: list[dict[str, Any]], headers: list[str]
) -> list[list[str]]:
    return [[str(row.get(header, "")) for header in headers] for row in row_dicts]


def add_slide_background(slide, prs, theme: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(
        *hex_to_rgb(theme["background_color"])
    )
    background.line.fill.background()


def add_top_accent(slide, theme: dict[str, Any]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Inches(12.0), Inches(0.12)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["accent_color"]))
    accent.line.fill.background()


def infer_template_name(slide_data: dict[str, Any]) -> str:
    explicit = str(slide_data.get("template", "")).strip().lower()
    if explicit:
        return explicit

    title = str(slide_data.get("title", "")).lower()
    slide_type = str(slide_data.get("type", "")).lower()

    if slide_type == "title":
        return "cover"
    if "经营概览" in title or "总览" in title:
        return "overview"
    if "核心指标" in title or "指标摘要" in title or "摘要" in title:
        return "kpi_summary"
    if "趋势" in title:
        return "weekly_trend"
    if "对比" in title or "分析" in title:
        return "comparison"
    if "结论" in title or "总结" in title or "summary" in title:
        return "conclusion"
    if "致谢" in title or "感谢" in title or "thank" in title:
        return "thanks"
    if "附录" in title:
        return "appendix"
    return ""


def split_lead_and_points(slide_data: dict[str, Any]) -> tuple[str, list[str]]:
    subtitle = str(slide_data.get("subtitle", "")).strip()
    bullets = [str(item).strip() for item in slide_data.get("bullets", []) if str(item).strip()]
    takeaways = [
        str(item).strip()
        for item in slide_data.get("takeaway_bullets", [])
        if str(item).strip()
    ]
    points = bullets or takeaways
    if subtitle:
        return subtitle, points
    if points:
        return points[0], points[1:]
    return "", []


def format_metric_value(value: Any) -> str:
    numeric = to_float(value)
    if numeric is None:
        return str(value)
    if abs(numeric - round(numeric)) < 1e-6:
        return f"{int(round(numeric)):,}"
    return f"{numeric:,.1f}"


def format_delta_text(current: float | None, previous: float | None) -> str:
    if current is None or previous is None:
        return ""
    if abs(previous) > 1e-9:
        ratio = (current - previous) / abs(previous)
        sign = "+" if ratio >= 0 else ""
        return f"{sign}{ratio * 100:.1f}% vs 上期"
    delta = current - previous
    sign = "+" if delta >= 0 else ""
    return f"{sign}{delta:,.0f} vs 上期"


def readable_metric_label(column: str) -> str:
    mapping = {
        "gmv": "GMV",
        "order_cnt": "订单量",
        "buyer_cnt": "买家数",
        "pay_buyer_cnt": "支付买家数",
        "arpu": "ARPU",
        "freq": "频次",
        "aov": "笔单价",
        "uv": "访客数",
        "pv": "浏览量",
        "ctr": "点击率",
        "conv_rate": "转化率",
        "member_gmv": "会员GMV",
        "member_order_cnt": "会员订单量",
        "member_buyer_cnt": "会员买家数",
    }
    normalized = column.strip().lower()
    return mapping.get(normalized, column.replace("_", " ").title())


def compact_context_label(value: Any) -> str:
    text = str(value).strip()
    if len(text) == 10 and text[4] == "-" and text[7] == "-":
        return text[5:]
    return text


def extract_metric_cards(
    query_data: dict[str, Any],
    max_metrics: int = 4,
) -> tuple[str, list[dict[str, str]]]:
    columns = query_data.get("columns", [])
    rows = query_data.get("rows", [])
    if not columns or not rows:
        return "", []

    latest = rows[-1]
    previous = rows[-2] if len(rows) > 1 else None
    context_column = next(
        (
            column
            for column in columns
            if any(to_float(row.get(column)) is None for row in rows)
        ),
        columns[0],
    )
    context_label = str(latest.get(context_column, ""))
    metric_columns = [
        column
        for column in columns
        if column != context_column
        and any(to_float(row.get(column)) is not None for row in rows)
    ][:max_metrics]

    cards = []
    for column in metric_columns:
        current = to_float(latest.get(column))
        previous_value = to_float(previous.get(column)) if previous else None
        cards.append(
            {
                "label": readable_metric_label(column),
                "value": format_metric_value(latest.get(column, "")),
                "delta": format_delta_text(current, previous_value),
            }
        )
    return context_label, cards


def add_card_shape(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str,
    line_color: str,
    radius_shape: str = "rounded",
):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius_shape == "rounded" else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(
        shape_type,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(*hex_to_rgb(fill_color))
    card.line.color.rgb = RGBColor(*hex_to_rgb(line_color))
    return card


def add_label_chip(slide, text: str, left: float, top: float, theme: dict[str, Any]) -> None:
    if not text:
        return
    chip = add_card_shape(
        slide,
        left,
        top,
        1.35 + min(len(text), 10) * 0.08,
        0.38,
        theme["soft_fill_color"],
        theme["soft_fill_color"],
    )
    box = add_textbox(slide, left + 0.14, top + 0.06, 1.8, 0.24)
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    apply_run_style(
        run,
        theme["font_family"],
        11,
        hex_to_rgb(theme["primary_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )
    chip.line.fill.background()


def add_hero_banner(
    slide,
    text: str,
    theme: dict[str, Any],
) -> None:
    """Full-width orange banner at the slide top with white bold text."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(0.72),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["primary_color"]))
    banner.line.fill.background()
    if text:
        banner_box = add_textbox(slide, 0.45, 0.14, 12.4, 0.44)
        banner_box.text_frame.word_wrap = True
        banner_run = banner_box.text_frame.paragraphs[0].add_run()
        banner_run.text = text
        apply_run_style(
            banner_run,
            theme["font_family"],
            theme["title_font_size"],
            (255, 255, 255),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )


def add_page_header(
    slide,
    title: str,
    theme: dict[str, Any],
    eyebrow: str = "经营分析",
    subtitle: str = "",
    banner_text: str = "",
) -> None:
    if banner_text:
        add_hero_banner(slide, banner_text, theme)
        # Compact header below banner: chip + title shifted down
        add_label_chip(slide, eyebrow, 0.75, 0.85, theme)
        title_box = add_textbox(slide, 0.75, 1.28, 7.6, 0.52)
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        paragraph = title_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = title
        apply_run_style(
            run,
            theme["font_family"],
            theme["title_font_size"] + 2,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        return
    add_label_chip(slide, eyebrow, 0.75, 0.55, theme)
    title_box = add_textbox(slide, 0.75, 0.98, 7.6, 0.72)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    paragraph = title_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    apply_run_style(
        run,
        theme["font_family"],
        theme["title_font_size"] + 2,
        hex_to_rgb(theme["text_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )
    if subtitle:
        subtitle_box = add_textbox(slide, 0.75, 1.56, 8.8, 0.4)
        subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
        subtitle_run.text = subtitle
        apply_run_style(
            subtitle_run,
            theme["font_family"],
            theme["body_font_size"] - 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )


def add_metric_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    metric: dict[str, str],
    theme: dict[str, Any],
    highlight: bool = False,
) -> None:
    fill_color = theme["soft_fill_color"] if highlight else theme["card_bg_color"]
    line_color = theme["warm_fill_color"] if highlight else theme["line_color"]
    add_card_shape(slide, left, top, width, height, fill_color, line_color)

    label_box = add_textbox(slide, left + 0.18, top + 0.18, width - 0.3, 0.26)
    label_run = label_box.text_frame.paragraphs[0].add_run()
    label_run.text = metric.get("label", "")
    apply_run_style(
        label_run,
        theme["font_family"],
        theme["body_font_size"] - 2,
        hex_to_rgb(theme["muted_text_color"]),
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )

    value_box = add_textbox(slide, left + 0.18, top + 0.48, width - 0.3, 0.5)
    value_run = value_box.text_frame.paragraphs[0].add_run()
    value_run.text = metric.get("value", "")
    apply_run_style(
        value_run,
        theme["font_family"],
        theme["hero_title_font_size"] + 4,
        hex_to_rgb(theme["text_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )

    delta = metric.get("delta", "")
    if delta:
        delta_box = add_textbox(slide, left + 0.18, top + height - 0.36, width - 0.3, 0.2)
        delta_run = delta_box.text_frame.paragraphs[0].add_run()
        delta_run.text = delta
        apply_run_style(
            delta_run,
            theme["font_family"],
            theme["footer_font_size"],
            hex_to_rgb(theme["primary_color"]),
            bold=highlight,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )


def resolve_chart_columns(
    query_data: dict[str, Any],
    category_column: str | None,
    value_columns: list[str] | None,
) -> tuple[str | None, list[str]]:
    columns = query_data.get("columns", [])
    rows = query_data.get("rows", [])
    if not columns or not rows:
        return None, []

    resolved_category = category_column or columns[0]
    resolved_values = list(value_columns or [])

    if not resolved_values:
        resolved_values = [
            column
            for column in columns
            if column != resolved_category
            and any(to_float(row.get(column)) is not None for row in rows)
        ][:3]

    return resolved_category, resolved_values


def build_chart_series(
    query_data: dict[str, Any],
    category_column: str,
    value_columns: list[str],
    max_points: int,
) -> tuple[list[str], list[tuple[str, list[float]]]]:
    rows = query_data.get("rows", [])[:max_points]
    categories: list[str] = []
    series_map = {column: [] for column in value_columns}

    for row in rows:
        categories.append(str(row.get(category_column, "")))
        for column in value_columns:
            value = to_float(row.get(column))
            series_map[column].append(0.0 if value is None else value)

    series = [(column, values) for column, values in series_map.items()]
    return categories, series


def chart_type_from_name(name: str):
    from pptx.enum.chart import XL_CHART_TYPE

    mapping = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    return mapping.get(str(name).lower(), XL_CHART_TYPE.LINE_MARKERS)


def infer_best_chart_type(template: str, explicit: str, num_categories: int) -> str:
    """Pick the best chart type based on template and data shape."""
    if explicit and explicit.lower() not in {"", "auto"}:
        return explicit.lower()
    recommendations = {
        "weekly_trend": "column",
        "comparison": "bar",
        "time_slot": "column",
        "price_band": "doughnut",
        "city_distribution": "bar",
        "aoi_distribution": "doughnut",
        "new_old_mix": "doughnut",
        "cluster_mix": "column",
    }
    recommended = recommendations.get(template, "")
    if recommended:
        return recommended
    if num_categories <= 5:
        return "column"
    return "bar"


def style_chart(chart, theme: dict[str, Any], chart_type_name: str) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_LEGEND_POSITION
    from pptx.util import Pt

    chart.has_title = False
    chart.has_legend = len(chart.series) > 1 and chart_type_name != "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        try:
            chart.legend.font.name = theme["cjk_font_family"] or theme["font_family"]
            chart.legend.font.size = Pt(theme["table_body_font_size"])
        except AttributeError:
            pass

    plot = chart.plots[0]
    plot.has_data_labels = False

    palette = theme["chart_palette"]
    for index, series in enumerate(chart.series):
        color = RGBColor(*hex_to_rgb(palette[index % len(palette)]))
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        except AttributeError:
            pass
        try:
            series.format.line.color.rgb = color
        except AttributeError:
            pass

    # Pie/doughnut charts have no axes – skip axis styling
    if chart_type_name not in {"pie", "doughnut"}:
        try:
            category_axis = chart.category_axis
            category_axis.tick_labels.font.name = theme["cjk_font_family"] or theme["font_family"]
            category_axis.tick_labels.font.size = Pt(theme["table_body_font_size"])
            category_axis.format.line.color.rgb = RGBColor(*hex_to_rgb(theme["line_color"]))
        except AttributeError:
            pass

        try:
            value_axis = chart.value_axis
            value_axis.has_major_gridlines = True
            value_axis.tick_labels.font.name = theme["cjk_font_family"] or theme["font_family"]
            value_axis.tick_labels.font.size = Pt(theme["table_body_font_size"])
            value_axis.format.line.color.rgb = RGBColor(
                *hex_to_rgb(theme["line_color"])
            )
            value_axis.major_gridlines.format.line.color.rgb = RGBColor(
                *hex_to_rgb(theme["line_color"])
            )
        except AttributeError:
            pass

    # Pie / Doughnut specific styling: show labels with percentages
    if chart_type_name in {"pie", "doughnut"}:
        try:
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_category_name = True
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.font.name = theme["cjk_font_family"] or theme["font_family"]
            data_labels.font.size = Pt(theme["table_body_font_size"] + 1)
        except AttributeError:
            pass


def add_chart_shape(
    slide,
    query_data: dict[str, Any],
    slide_data: dict[str, Any],
    theme: dict[str, Any],
    left: float = 0.7,
    top: float = 2.2,
    width: float = 12.0,
    height: float = 4.25,
):
    from pptx.chart.data import ChartData
    from pptx.util import Inches

    category_column, value_columns = resolve_chart_columns(
        query_data,
        slide_data.get("category_column"),
        slide_data.get("value_columns"),
    )
    if not category_column or not value_columns:
        return None

    max_points = int(slide_data.get("max_points", 12))
    categories, series_data = build_chart_series(
        query_data, category_column, value_columns, max_points
    )
    if not categories or not series_data:
        return None

    template = infer_template_name(slide_data)
    explicit_chart_type = str(slide_data.get("chart_type", ""))
    chart_type_name = infer_best_chart_type(template, explicit_chart_type, len(categories))

    # Doughnut/pie charts work best with a single value series
    if chart_type_name in {"pie", "doughnut"} and len(series_data) > 1:
        series_data = series_data[:1]

    chart_data = ChartData()
    chart_data.categories = categories
    for series_name, values in series_data:
        chart_data.add_series(series_name, values)

    chart_shape = slide.shapes.add_chart(
        chart_type_from_name(chart_type_name),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    )
    chart = chart_shape.chart
    style_chart(chart, theme, chart_type_name)
    return chart_shape


def render_title_slide(prs, slide_data: dict[str, Any], theme: dict[str, Any]):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, theme)
    add_logo_to_slide(slide, theme, 0.8, 0.15, 0.38)
    add_label_chip(slide, "品牌经营分析", 0.8, 0.65, theme)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.84), Inches(1.35), Inches(0.08), Inches(4.95)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["primary_color"]))
    accent.line.fill.background()

    title_box = add_textbox(slide, 1.18, 1.2, 6.4, 2.3)
    title_box.text_frame.word_wrap = True
    paragraph = title_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = slide_data.get("title", "")
    apply_run_style(
        run,
        theme["font_family"],
        theme["hero_title_font_size"] + 10,
        hex_to_rgb(theme["text_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )

    deck_label = add_textbox(slide, 1.18, 3.85, 4.3, 0.22)
    deck_run = deck_label.text_frame.paragraphs[0].add_run()
    deck_run.text = "淘宝闪购经营分析作战室"
    apply_run_style(
        deck_run,
        theme["font_family"],
        theme["footer_font_size"] + 2,
        hex_to_rgb(theme["primary_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.12), Inches(4.95), Inches(0.03)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["line_color"]))
    divider.line.fill.background()

    subtitle_box = add_textbox(slide, 1.18, 4.28, 5.95, 0.38)
    subtitle_box.text_frame.word_wrap = True
    subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
    subtitle_run.text = slide_data.get("subtitle", "")
    apply_run_style(
        subtitle_run,
        theme["font_family"],
        theme["subtitle_font_size"] - 1,
        hex_to_rgb(theme["muted_text_color"]),
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )

    subline_box = add_textbox(slide, 1.18, 4.78, 6.2, 0.26)
    subline_run = subline_box.text_frame.paragraphs[0].add_run()
    subline_run.text = "品牌经营分析 / 关键指标摘要 / 重点专题拆解"
    apply_run_style(
        subline_run,
        theme["font_family"],
        theme["footer_font_size"] + 1,
        hex_to_rgb(theme["muted_text_color"]),
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )

    info_panel = add_card_shape(
        slide,
        8.78,
        1.02,
        2.86,
        5.55,
        theme["card_bg_color"],
        theme["line_color"],
        radius_shape="square",
    )
    info_panel.line.color.rgb = RGBColor(*hex_to_rgb(theme["line_color"]))
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11.36), Inches(1.02), Inches(0.12), Inches(5.55)
    )
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["primary_color"]))
    side_bar.line.fill.background()

    cover_kicker = add_textbox(slide, 9.02, 1.28, 2.0, 0.2)
    cover_run = cover_kicker.text_frame.paragraphs[0].add_run()
    cover_run.text = "MONTHLY REVIEW"
    apply_run_style(
        cover_run,
        theme["font_family"],
        theme["footer_font_size"] + 1,
        hex_to_rgb(theme["primary_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )

    info_rows = [
        ("报告主题", slide_data.get("title", "")),
        ("观察周期", slide_data.get("subtitle", "")),
        ("材料类型", "品牌经营分析报告"),
    ]
    row_top = 1.88
    for label, value in info_rows:
        label_box = add_textbox(slide, 9.02, row_top, 1.0, 0.18)
        label_run = label_box.text_frame.paragraphs[0].add_run()
        label_run.text = label
        apply_run_style(
            label_run,
            theme["font_family"],
            theme["footer_font_size"] + 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        value_box = add_textbox(slide, 9.02, row_top + 0.24, 1.95, 0.6)
        value_box.text_frame.word_wrap = True
        value_run = value_box.text_frame.paragraphs[0].add_run()
        value_run.text = str(value)
        apply_run_style(
            value_run,
            theme["font_family"],
            theme["body_font_size"] - 1,
            hex_to_rgb(theme["text_color"]),
            bold=True if label != "观察周期" else False,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(9.02), Inches(row_top + 0.94), Inches(2.0), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["line_color"]))
        line.line.fill.background()
        row_top += 1.1

    insight_card = add_card_shape(
        slide,
        9.02,
        5.28,
        1.98,
        0.9,
        theme["soft_fill_color"],
        theme["soft_fill_color"],
        radius_shape="square",
    )
    insight_card.line.color.rgb = RGBColor(*hex_to_rgb(theme["soft_fill_color"]))
    insight_tag = add_textbox(slide, 9.18, 5.42, 0.8, 0.16)
    insight_tag_run = insight_tag.text_frame.paragraphs[0].add_run()
    insight_tag_run.text = "摘要提示"
    apply_run_style(
        insight_tag_run,
        theme["font_family"],
        theme["footer_font_size"],
        hex_to_rgb(theme["primary_color"]),
        bold=True,
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )
    insight_box = add_textbox(slide, 9.18, 5.66, 1.6, 0.28)
    insight_box.text_frame.word_wrap = True
    insight_run = insight_box.text_frame.paragraphs[0].add_run()
    insight_run.text = "先看品牌整体经营，再看结构与参考对比。"
    apply_run_style(
        insight_run,
        theme["font_family"],
        theme["footer_font_size"] + 1,
        hex_to_rgb(theme["text_color"]),
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )
    meta_box = add_textbox(slide, 1.18, 6.15, 6.8, 0.24)
    meta_run = meta_box.text_frame.paragraphs[0].add_run()
    meta_run.text = "围绕本期经营表现、关键指标变化和重点结构问题展开。"
    apply_run_style(
        meta_run,
        theme["font_family"],
        theme["footer_font_size"] + 1,
        hex_to_rgb(theme["muted_text_color"]),
        cjk_font_name=theme["cjk_font_family"],
        language=theme["language"],
    )


def render_schema_slide(
    prs, slide_data: dict[str, Any], payload: dict[str, Any], theme: dict[str, Any]
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, theme)
    add_page_header(slide, slide_data["title"], theme, eyebrow="附录")

    intro_text = slide_data.get("intro_text", "")
    if intro_text:
        add_bullets_box(slide, [intro_text], 0.75, 1.75, 12.0, 0.45, theme)

    tables = slide_data.get("tables", [])
    if not tables:
        add_bullets_box(
            slide,
            ["这一页没有配置 schema 表。"],
            0.7,
            1.5,
            12.0,
            1.0,
            theme,
        )
        return

    available_height = 4.7
    per_table_height = available_height / len(tables)
    top = 2.1
    for table_name in tables:
        schema_rows = payload["tables"].get(table_name, [])
        headers = ["column_name", "data_type", "is_partition_key", "column_comment"]
        normalized = normalize_rows(schema_rows, headers)[:8]
        if not normalized:
            normalized = [["没有查到 schema 结果", "", "", ""]]
        add_table_shape(
            slide,
            headers,
            normalized,
            0.7,
            top,
            12.0,
            per_table_height - 0.15,
            theme,
        )
        top += per_table_height

    add_footer(slide, slide_data.get("speaker_notes", ""), theme)


def render_table_slide(
    prs, slide_data: dict[str, Any], payload: dict[str, Any], theme: dict[str, Any]
):
    template = infer_template_name(slide_data)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, theme)

    query_name = slide_data["query"]
    query_data = payload["queries"].get(query_name)
    early_takeaways = [
        str(item).strip()
        for item in slide_data.get("takeaway_bullets", [])
        if str(item).strip()
    ]
    if template != "overview":
        eyebrow = "数据页"
        if template in {"kpi_summary", "time_slot", "price_band", "time_price_matrix", "city_distribution", "aoi_distribution", "new_old_mix", "cluster_mix"}:
            eyebrow = "结构分析"
        elif template == "appendix":
            eyebrow = "附录"
        banner_text = ""
        if template not in {"kpi_summary", "appendix"}:
            banner_text = early_takeaways[0] if early_takeaways else slide_data["title"]
        add_page_header(
            slide,
            slide_data["title"],
            theme,
            eyebrow=eyebrow,
            banner_text=banner_text,
        )
    if not query_data:
        add_bullets_box(
            slide,
            [f"没有在 payload 里找到查询结果：{query_name}"],
            0.7,
            2.1,
            12.0,
            1.0,
            theme,
        )
        return

    headers = query_data["columns"]
    max_rows = int(slide_data.get("max_rows", 12))
    rows = normalize_rows(query_data["rows"][:max_rows], headers)
    if not rows:
        rows = [["没有返回数据"] + [""] * (len(headers) - 1)]
    context_label, metric_cards = extract_metric_cards(query_data, 6)
    takeaways = [
        str(item).strip()
        for item in slide_data.get("takeaway_bullets", [])
        if str(item).strip()
    ]

    if template == "overview":
        latest_row = query_data["rows"][-1]
        context_text = context_label or latest_row.get(headers[0], "")
        hero_metric = metric_cards[0] if metric_cards else {"label": "核心指标", "value": "-", "delta": ""}
        support_metrics = metric_cards[1:4]
        lead_text = takeaways[0] if takeaways else "本期整体经营表现继续走强，核心指标延续上升。"
        support_text = takeaways[1] if len(takeaways) > 1 else "经营盘面保持改善，末端走势仍在抬升。"
        support_query_name = str(slide_data.get("support_query", "")).strip()
        support_query_data = payload["queries"].get(support_query_name, {})
        support_headers = support_query_data.get("columns", [])
        support_rows = support_query_data.get("rows", [])

        add_label_chip(slide, "经营概览", 0.78, 0.62, theme)
        title_box = add_textbox(slide, 0.78, 1.02, 5.8, 0.72)
        title_run = title_box.text_frame.paragraphs[0].add_run()
        title_run.text = slide_data["title"]
        apply_run_style(
            title_run,
            theme["font_family"],
            theme["title_font_size"] + 4,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        period_box = add_textbox(slide, 9.25, 0.9, 2.6, 0.25)
        period_run = period_box.text_frame.paragraphs[0].add_run()
        period_run.text = f"观察窗口  {context_text}"
        apply_run_style(
            period_run,
            theme["font_family"],
            theme["footer_font_size"] + 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        hero_card = add_card_shape(
            slide,
            0.78,
            1.95,
            6.55,
            3.45,
            theme["soft_fill_color"],
            theme["warm_fill_color"],
        )
        hero_card.line.fill.background()
        add_label_chip(slide, "本期判断", 1.0, 2.18, theme)

        lead_box = add_textbox(slide, 1.0, 2.62, 5.65, 0.9)
        lead_box.text_frame.word_wrap = True
        lead_run = lead_box.text_frame.paragraphs[0].add_run()
        lead_run.text = lead_text
        apply_run_style(
            lead_run,
            theme["font_family"],
            theme["subtitle_font_size"] + 4,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        hero_label = add_textbox(slide, 1.0, 3.75, 1.5, 0.22)
        hero_label_run = hero_label.text_frame.paragraphs[0].add_run()
        hero_label_run.text = hero_metric.get("label", "")
        apply_run_style(
            hero_label_run,
            theme["font_family"],
            theme["body_font_size"] - 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        hero_value = add_textbox(slide, 1.0, 4.0, 3.3, 0.75)
        hero_value_run = hero_value.text_frame.paragraphs[0].add_run()
        hero_value_run.text = hero_metric.get("value", "-")
        apply_run_style(
            hero_value_run,
            theme["font_family"],
            theme["hero_title_font_size"] + 14,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        hero_delta = add_textbox(slide, 4.18, 4.18, 1.8, 0.24)
        hero_delta_run = hero_delta.text_frame.paragraphs[0].add_run()
        hero_delta_run.text = hero_metric.get("delta", "")
        apply_run_style(
            hero_delta_run,
            theme["font_family"],
            theme["body_font_size"] - 1,
            hex_to_rgb(theme["primary_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        support_box = add_textbox(slide, 1.0, 4.82, 5.65, 0.34)
        support_box.text_frame.word_wrap = True
        support_run = support_box.text_frame.paragraphs[0].add_run()
        support_run.text = support_text
        apply_run_style(
            support_run,
            theme["font_family"],
            theme["body_font_size"] - 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        for index, metric in enumerate(support_metrics[:3]):
            card = add_card_shape(
                slide,
                7.62,
                1.95 + index * 1.18,
                4.63,
                1.0,
                theme["card_bg_color"],
                theme["line_color"],
            )
            card.line.fill.background()
            label_box = add_textbox(slide, 7.9, 2.18 + index * 1.18, 1.7, 0.2)
            label_run = label_box.text_frame.paragraphs[0].add_run()
            label_run.text = metric.get("label", "")
            apply_run_style(
                label_run,
                theme["font_family"],
                theme["footer_font_size"] + 1,
                hex_to_rgb(theme["muted_text_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
            value_box = add_textbox(slide, 7.9, 2.43 + index * 1.18, 2.0, 0.34)
            value_run = value_box.text_frame.paragraphs[0].add_run()
            value_run.text = metric.get("value", "-")
            apply_run_style(
                value_run,
                theme["font_family"],
                theme["subtitle_font_size"] + 4,
                hex_to_rgb(theme["text_color"]),
                bold=True,
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
            delta_box = add_textbox(slide, 10.1, 2.48 + index * 1.18, 1.8, 0.2)
            delta_run = delta_box.text_frame.paragraphs[0].add_run()
            delta_run.text = metric.get("delta", "")
            apply_run_style(
                delta_run,
                theme["font_family"],
                theme["footer_font_size"] + 1,
                hex_to_rgb(theme["primary_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )

        if support_rows and support_headers:
            strip_rows = support_rows[-5:] if len(support_rows) >= 5 else support_rows
            strip_title_text = "月内周趋势"
            strip_category = support_headers[0]
            strip_value_column = next(
                (
                    column
                    for column in support_headers[1:]
                    if any(to_float(row.get(column)) is not None for row in support_rows)
                ),
                support_headers[1] if len(support_headers) > 1 else support_headers[0],
            )
        else:
            strip_rows = query_data["rows"][-6:] if len(query_data["rows"]) >= 6 else query_data["rows"]
            strip_title_text = "近六日经营走势"
            strip_category = headers[0]
            strip_value_column = headers[1] if len(headers) > 1 else headers[0]

        add_card_shape(
            slide,
            0.78,
            5.72,
            11.48,
            1.0,
            theme["card_bg_color"],
            theme["line_color"],
        )
        strip_title = add_textbox(slide, 1.0, 5.92, 1.8, 0.18)
        strip_run = strip_title.text_frame.paragraphs[0].add_run()
        strip_run.text = strip_title_text
        apply_run_style(
            strip_run,
            theme["font_family"],
            theme["footer_font_size"] + 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        card_width = 1.55
        start_left = 2.8
        for index, row in enumerate(strip_rows[:6]):
            tile = add_card_shape(
                slide,
                start_left + index * 1.42,
                5.88,
                card_width,
                0.58,
                theme["background_color"],
                theme["line_color"],
            )
            tile.line.fill.background()
            day_box = add_textbox(slide, start_left + 0.12 + index * 1.42, 6.0, 0.52, 0.14)
            day_run = day_box.text_frame.paragraphs[0].add_run()
            day_run.text = compact_context_label(row.get(strip_category, ""))
            apply_run_style(
                day_run,
                theme["font_family"],
                theme["footer_font_size"],
                hex_to_rgb(theme["muted_text_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
            value_box = add_textbox(slide, start_left + 0.12 + index * 1.42, 6.18, 1.2, 0.18)
            value_run = value_box.text_frame.paragraphs[0].add_run()
            value_run.text = format_metric_value(row.get(strip_value_column, ""))
            apply_run_style(
                value_run,
                theme["font_family"],
                theme["body_font_size"] - 1,
                hex_to_rgb(theme["text_color"]),
                bold=True,
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
    elif template == "kpi_summary":
        lead_text = takeaways[0] if takeaways else "先看本期最值得关注的核心指标。"
        add_card_shape(
            slide,
            0.75,
            2.02,
            11.8,
            0.72,
            theme["soft_fill_color"],
            theme["soft_fill_color"],
        )
        lead_box = add_textbox(slide, 0.95, 2.2, 11.0, 0.28)
        lead_run = lead_box.text_frame.paragraphs[0].add_run()
        lead_run.text = lead_text
        apply_run_style(
            lead_run,
            theme["font_family"],
            theme["body_font_size"],
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        for index, metric in enumerate(metric_cards[:6]):
            row = index // 3
            col = index % 3
            add_metric_card(
                slide,
                0.75 + col * 3.95,
                3.0 + row * 1.72,
                3.55,
                1.45,
                metric,
                theme,
                highlight=index == 0,
            )

        if context_label:
            foot_box = add_textbox(slide, 0.95, 6.55, 6.0, 0.24)
            foot_run = foot_box.text_frame.paragraphs[0].add_run()
            foot_run.text = f"指标口径参考 | 最新数据点：{context_label}"
            apply_run_style(
                foot_run,
                theme["font_family"],
                theme["footer_font_size"],
                hex_to_rgb(theme["muted_text_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
    else:
        if takeaways:
            add_card_shape(
                slide,
                0.75,
                2.0,
                11.8,
                0.78,
                theme["soft_fill_color"],
                theme["soft_fill_color"],
            )
            add_bullets_box(slide, takeaways[:2], 0.98, 2.18, 11.0, 0.36, theme, font_size=13)
        add_card_shape(
            slide,
            0.75,
            2.95,
            11.8,
            3.55,
            theme["card_bg_color"],
            theme["line_color"],
        )
        add_table_shape(slide, headers, rows, 0.95, 3.2, 11.4, 3.05, theme)

    add_footer(slide, slide_data.get("speaker_notes", ""), theme)


def render_chart_slide(
    prs, slide_data: dict[str, Any], payload: dict[str, Any], theme: dict[str, Any]
):
    template = infer_template_name(slide_data)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, theme)

    query_name = slide_data["query"]
    query_data = payload["queries"].get(query_name)
    early_takeaways = [
        str(item).strip()
        for item in slide_data.get("takeaway_bullets", [])
        if str(item).strip()
    ]
    if template != "comparison":
        eyebrow = "分析页"
        if template == "weekly_trend":
            eyebrow = "本月趋势"
        elif template in {"time_slot", "price_band", "city_distribution", "aoi_distribution", "new_old_mix"}:
            eyebrow = "结构分析"
        banner_text = early_takeaways[0] if early_takeaways else slide_data["title"]
        add_page_header(
            slide,
            slide_data["title"],
            theme,
            eyebrow=eyebrow,
            banner_text=banner_text,
        )

    if not query_data:
        add_bullets_box(
            slide,
            [f"没有在 payload 里找到查询结果：{query_name}"],
            0.7,
            2.0,
            12.0,
            1.0,
            theme,
        )
        return

    takeaways = [
        str(item).strip()
        for item in slide_data.get("takeaway_bullets", [])
        if str(item).strip()
    ]

    chart_left = 0.95
    chart_top = 2.2
    chart_width = 11.4
    chart_height = 4.1
    if template == "comparison":
        lead_text = takeaways[0] if takeaways else "先看趋势主线，再把参考对比当作辅助说明。"
        support_points = takeaways[1:3]
        context_label, metric_cards = extract_metric_cards(query_data, 2)
        opinion_text = support_points[0] if support_points else "最新一期 GMV 延续上行，当前处于稳步抬升区间。"

        add_label_chip(slide, "分析对比", 0.78, 0.62, theme)
        title_box = add_textbox(slide, 0.78, 1.02, 6.0, 0.72)
        title_run = title_box.text_frame.paragraphs[0].add_run()
        title_run.text = slide_data["title"]
        apply_run_style(
            title_run,
            theme["font_family"],
            theme["title_font_size"] + 4,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        intro_box = add_textbox(slide, 0.8, 1.56, 6.8, 0.22)
        intro_run = intro_box.text_frame.paragraphs[0].add_run()
        intro_run.text = lead_text
        apply_run_style(
            intro_run,
            theme["font_family"],
            theme["body_font_size"] - 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        opinion_box = add_textbox(slide, 0.8, 1.98, 8.0, 0.46)
        opinion_box.text_frame.word_wrap = True
        opinion_run = opinion_box.text_frame.paragraphs[0].add_run()
        opinion_run.text = opinion_text
        apply_run_style(
            opinion_run,
            theme["font_family"],
            theme["subtitle_font_size"] + 4,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        chart_card = add_card_shape(
            slide,
            0.78,
            2.72,
            8.48,
            3.95,
            theme["card_bg_color"],
            theme["line_color"],
        )
        chart_card.line.fill.background()
        add_label_chip(slide, "品牌主视角", 1.02, 2.32, theme)
        chart_lead = add_textbox(slide, 1.02, 3.08, 7.35, 0.26)
        chart_lead_run = chart_lead.text_frame.paragraphs[0].add_run()
        chart_lead_run.text = "趋势图只负责支撑观点，参考关系只做轻量标注。"
        apply_run_style(
            chart_lead_run,
            theme["font_family"],
            theme["footer_font_size"] + 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        evidence_card = add_card_shape(
            slide,
            9.58,
            3.05,
            2.58,
            3.62,
            theme["soft_fill_color"],
            theme["warm_fill_color"],
        )
        evidence_card.line.fill.background()
        add_label_chip(slide, "证据摘要", 9.82, 3.28, theme)
        rail_title = add_textbox(slide, 9.82, 3.72, 1.95, 0.72)
        rail_title.text_frame.word_wrap = True
        rail_run = rail_title.text_frame.paragraphs[0].add_run()
        rail_run.text = "大盘仅作辅助参考，不与品牌主视角并列。"
        apply_run_style(
            rail_run,
            theme["font_family"],
            theme["body_font_size"] + 1,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        evidence_points = support_points[:2] or ["最新一周 GMV 延续上行。", "末端点位仍在抬升。"]
        for index, bullet in enumerate(evidence_points):
            bullet_box = add_textbox(slide, 9.84, 4.56 + index * 0.54, 1.95, 0.34)
            bullet_box.text_frame.word_wrap = True
            bullet_run = bullet_box.text_frame.paragraphs[0].add_run()
            bullet_run.text = f"{index + 1:02d}  {bullet}"
            apply_run_style(
                bullet_run,
                theme["font_family"],
                theme["footer_font_size"] + 2,
                hex_to_rgb(theme["text_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )

        if metric_cards:
            add_metric_card(slide, 9.82, 5.78, 1.95, 0.98, metric_cards[0], theme, highlight=True)
        if context_label:
            ref_box = add_textbox(slide, 9.84, 6.48, 1.8, 0.18)
            ref_run = ref_box.text_frame.paragraphs[0].add_run()
            ref_run.text = f"参考窗口  {context_label}"
            apply_run_style(
                ref_run,
                theme["font_family"],
                theme["footer_font_size"],
                hex_to_rgb(theme["muted_text_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )

        chart_left = 1.02
        chart_top = 3.48
        chart_width = 8.04
        chart_height = 2.65
    elif takeaways:
        add_card_shape(
            slide,
            0.75,
            2.03,
            11.8,
            0.78,
            theme["soft_fill_color"],
            theme["soft_fill_color"],
        )
        add_bullets_box(slide, takeaways[:2], 0.98, 2.18, 11.0, 0.36, theme, font_size=13)
        add_card_shape(
            slide,
            0.75,
            3.0,
            11.8,
            3.8,
            theme["card_bg_color"],
            theme["line_color"],
        )
        chart_left = 0.95
        chart_top = 3.25
        chart_width = 11.4
        chart_height = 3.3

    chart_shape = add_chart_shape(
        slide,
        query_data,
        slide_data,
        theme,
        left=chart_left,
        top=chart_top,
        width=chart_width,
        height=chart_height,
    )
    if chart_shape is None:
        add_bullets_box(
            slide,
            ["这一页没有足够的维度列或数值列，暂时没法生成原生图表。"],
            0.7,
            2.0,
            12.0,
            1.0,
            theme,
        )
        return

    add_footer(slide, slide_data.get("speaker_notes", ""), theme)


def render_bullet_slide(prs, slide_data: dict[str, Any], theme: dict[str, Any]):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    template = infer_template_name(slide_data)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background(slide, prs, theme)
    lead_text, points = split_lead_and_points(slide_data)

    if template == "thanks":
        add_label_chip(slide, "报告结束", 0.82, 0.72, theme)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(1.42), Inches(0.12), Inches(2.72)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["primary_color"]))
        bar.line.fill.background()

        title_box = add_textbox(slide, 1.22, 1.46, 5.0, 0.95)
        title_run = title_box.text_frame.paragraphs[0].add_run()
        title_run.text = slide_data["title"]
        apply_run_style(
            title_run,
            theme["font_family"],
            theme["hero_title_font_size"] + 2,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        subtitle_text = lead_text or "欢迎继续沟通具体问题。"
        subtitle_box = add_textbox(slide, 1.22, 2.46, 5.2, 0.42)
        subtitle_run = subtitle_box.text_frame.paragraphs[0].add_run()
        subtitle_run.text = subtitle_text
        apply_run_style(
            subtitle_run,
            theme["font_family"],
            theme["body_font_size"],
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        closing_note = add_textbox(slide, 1.22, 3.28, 4.8, 0.4)
        closing_run = closing_note.text_frame.paragraphs[0].add_run()
        closing_run.text = "如需继续展开时段、价格带或城市结构，可在此基础上继续拆解。"
        apply_run_style(
            closing_run,
            theme["font_family"],
            theme["footer_font_size"] + 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        info_card = add_card_shape(
            slide,
            8.28,
            1.52,
            3.08,
            2.56,
            theme["card_bg_color"],
            theme["line_color"],
            radius_shape="square",
        )
        info_card.line.color.rgb = RGBColor(*hex_to_rgb(theme["line_color"]))
        info_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.28), Inches(1.52), Inches(0.12), Inches(2.56)
        )
        info_bar.fill.solid()
        info_bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["primary_color"]))
        info_bar.line.fill.background()

        detail_rows = [
            ("报告来源", "品牌经营分析作战室"),
            ("输出形式", "可编辑经营分析PPT"),
            ("下一步", subtitle_text),
        ]
        detail_top = 1.78
        for label, value in detail_rows:
            label_box = add_textbox(slide, 8.62, detail_top, 0.9, 0.18)
            label_run = label_box.text_frame.paragraphs[0].add_run()
            label_run.text = label
            apply_run_style(
                label_run,
                theme["font_family"],
                theme["footer_font_size"] + 1,
                hex_to_rgb(theme["muted_text_color"]),
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
            value_box = add_textbox(slide, 8.62, detail_top + 0.24, 2.15, 0.4)
            value_box.text_frame.word_wrap = True
            value_run = value_box.text_frame.paragraphs[0].add_run()
            value_run.text = value
            apply_run_style(
                value_run,
                theme["font_family"],
                theme["footer_font_size"] + 2,
                hex_to_rgb(theme["text_color"]),
                bold=label != "下一步",
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
            if label != "下一步":
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(8.62), Inches(detail_top + 0.82), Inches(2.15), Inches(0.02)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["line_color"]))
                line.line.fill.background()
            detail_top += 0.86
        return

    if template != "conclusion":
        add_page_header(
            slide,
            slide_data["title"],
            theme,
            eyebrow="总结页" if template == "conclusion" else "结论页",
        )

    if template == "conclusion":
        add_label_chip(slide, "本期结论", 0.78, 0.62, theme)
        title_box = add_textbox(slide, 0.78, 1.02, 5.6, 0.72)
        title_run = title_box.text_frame.paragraphs[0].add_run()
        title_run.text = slide_data["title"]
        apply_run_style(
            title_run,
            theme["font_family"],
            theme["title_font_size"] + 4,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )
        lead_statement = lead_text or (points[0] if points else "本期经营延续改善，后续重点看结构优化和延续性。")
        roadmap_points = points or [str(item) for item in slide_data.get("bullets", [])]
        if lead_text and roadmap_points and roadmap_points[0] == lead_text:
            roadmap_points = roadmap_points[1:]

        lead_card = add_card_shape(
            slide,
            0.78,
            2.0,
            5.35,
            4.62,
            theme["soft_fill_color"],
            theme["warm_fill_color"],
        )
        lead_card.line.fill.background()

        flag = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.98), Inches(2.25), Inches(0.12), Inches(1.02)
        )
        flag.fill.solid()
        flag.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["primary_color"]))
        flag.line.fill.background()

        lead_label = add_textbox(slide, 1.22, 2.24, 1.5, 0.22)
        lead_label_run = lead_label.text_frame.paragraphs[0].add_run()
        lead_label_run.text = "总判断"
        apply_run_style(
            lead_label_run,
            theme["font_family"],
            theme["footer_font_size"] + 2,
            hex_to_rgb(theme["primary_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        lead_box = add_textbox(slide, 1.22, 2.68, 4.45, 1.65)
        lead_box.text_frame.word_wrap = True
        lead_run = lead_box.text_frame.paragraphs[0].add_run()
        lead_run.text = lead_statement
        apply_run_style(
            lead_run,
            theme["font_family"],
            theme["subtitle_font_size"] + 8,
            hex_to_rgb(theme["text_color"]),
            bold=True,
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        note_box = add_textbox(slide, 1.22, 4.72, 4.3, 0.7)
        note_box.text_frame.word_wrap = True
        note_run = note_box.text_frame.paragraphs[0].add_run()
        note_run.text = "这一页先给结论，再把后续关注点拆成三条，方便管理层直接抓重点。"
        apply_run_style(
            note_run,
            theme["font_family"],
            theme["body_font_size"] - 1,
            hex_to_rgb(theme["muted_text_color"]),
            cjk_font_name=theme["cjk_font_family"],
            language=theme["language"],
        )

        rail = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(7.28), Inches(2.25), Inches(0.03), Inches(3.72)
        )
        rail.fill.solid()
        rail.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["warm_fill_color"]))
        rail.line.fill.background()

        for index, point in enumerate(roadmap_points[:3]):
            center_y = 2.55 + index * 1.2
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(6.98), Inches(center_y), Inches(0.46), Inches(0.46)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(*hex_to_rgb(theme["background_color"]))
            marker.line.color.rgb = RGBColor(*hex_to_rgb(theme["primary_color"]))

            marker_text = add_textbox(slide, 7.02, center_y + 0.11, 0.36, 0.14)
            marker_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            marker_run = marker_text.text_frame.paragraphs[0].add_run()
            marker_run.text = f"{index + 1:02d}"
            apply_run_style(
                marker_run,
                theme["font_family"],
                theme["footer_font_size"],
                hex_to_rgb(theme["primary_color"]),
                bold=True,
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )

            point_box = add_textbox(slide, 7.58, center_y - 0.04, 4.2, 0.65)
            point_box.text_frame.word_wrap = True
            point_run = point_box.text_frame.paragraphs[0].add_run()
            point_run.text = str(point)
            apply_run_style(
                point_run,
                theme["font_family"],
                theme["body_font_size"] + 1,
                hex_to_rgb(theme["text_color"]),
                bold=index == 0,
                cjk_font_name=theme["cjk_font_family"],
                language=theme["language"],
            )
    else:
        if lead_text:
            add_bullets_box(
                slide,
                [lead_text],
                0.75,
                1.85,
                11.8,
                0.45,
                theme,
                font_size=theme["subtitle_font_size"],
            )
        add_bullets_box(slide, points or slide_data.get("bullets", []), 0.9, 2.5, 11.3, 4.1, theme)

    add_footer(slide, slide_data.get("speaker_notes", ""), theme)


def main() -> None:
    ensure_dependencies()
    from pptx import Presentation
    from pptx.util import Inches

    args = parse_args()
    config_path = Path(args.config).resolve()
    payload_path = Path(args.payload).resolve()
    plan_path = Path(args.plan).resolve()

    base_config = load_yaml(config_path)
    payload = load_json(payload_path)
    context = payload["metadata"].get("variables", {})
    rendered_config = deep_render(base_config, context)
    plan = load_json(plan_path)

    theme = build_theme(rendered_config.get("report", {}).get("theme", {}))
    resolved_logo = resolve_logo_path(config_path, theme)
    if resolved_logo:
        theme["_logo_resolved_path"] = resolved_logo

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in plan.get("slides", []):
        slide_type = slide_data.get("type")
        if slide_type == "title":
            render_title_slide(prs, slide_data, theme)
        elif slide_type == "schema":
            render_schema_slide(prs, slide_data, payload, theme)
        elif slide_type == "table":
            render_table_slide(prs, slide_data, payload, theme)
        elif slide_type == "chart":
            render_chart_slide(prs, slide_data, payload, theme)
        elif slide_type == "ai_bullets":
            render_bullet_slide(prs, slide_data, theme)
        else:
            raise ValueError(f"Unsupported slide type: {slide_type}")

    if args.output:
        output_path = Path(args.output).resolve()
    else:
        report = rendered_config.get("report", {})
        default_output = report.get("output_filename", "report.pptx")
        output_path = output_dir_from_config(config_path, rendered_config) / default_output

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"ppt={output_path}")


if __name__ == "__main__":
    main()
